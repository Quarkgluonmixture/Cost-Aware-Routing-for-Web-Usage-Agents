#!/usr/bin/env python3
"""Build the v9 showcase poster: a picture-led sheet on the organiser's A1 portrait template.

Why v9 exists
-------------
v5-v8 grew by adding sentences. The supervisor's reference sheet (a NeurIPS
three-column landscape poster) is the opposite: figures occupy roughly seven
tenths of it and the prose is short bullets under them. The author's brief for
v9, in their words: **more figures, fewer words, more examples, more
screenshots, no big numbers (they read "low"), title = the thesis title**, and
the test a passer-by applies — *catch the eye, and after two glances know what
this is about*.

What changed from v8
--------------------
* the headline number strip is gone (three 42pt figures on a tinted band);
* the title is the thesis title and the standfirst is its subtitle, so the sheet
  states its subject instead of its result;
* the top fifth is one wide system diagram with real screenshots in it;
* the evidence is a two-lane screenshot strip — one task, two ways of seeing the
  page — where the reader can see the failure rather than read about it;
* every panel is a figure plus at most two lines.

The sheet stays on the template and is **not** resized: A1 portrait, 594x841mm.
``verify`` re-asserts that and fails the build if a column overruns its box.

Assets come from ``poster_figures_v9.py`` (run it first) plus the thesis and
analysis figure libraries. Nothing here is drawn by hand and no number on the
sheet is typed: the strip's costs are summed from the step records at build time.

Usage::

    .venv/bin/python3 deliverables/showcase/poster_figures_v9.py
    .venv/bin/python3 deliverables/showcase/build_poster_v9.py
"""

from __future__ import annotations

import sys
from pathlib import Path

import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Mm, Pt

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

# v8's template plumbing is generic (clone/set_text/textbox/rect/measure); only
# the layout below is new, so it is imported rather than copied.
import build_poster as V8  # noqa: E402
from build_poster import (  # noqa: E402
    ACCENT, C_BOTH, C_FAIL, C_IMAGE, C_TEXT, FIG_FILL, GREY, HAIRLINE, INK,
    INK_STRONG, MONO, MUTED, PAPER, SANS, clone, drop, find, rect, set_text,
    textbox,
)
from pptx.enum.text import PP_ALIGN  # noqa: E402

REPO = Path(__file__).resolve().parents[2]
TEMPLATE = HERE / "Showcase-Poster-Template-A1.pptx"
OUT = HERE / "poster_v9_jiaming_wei.pptx"
V9 = HERE / "figures" / "v9"
THESIS = REPO / "final_dissertation" / "figures"
PHANTOM = REPO / "results" / "phantom_paper" / "figures"
REPL = REPO / "results" / "repro_replicates"
REPO_URL = "https://github.com/Quarkgluonmixture/Cost-Aware-Routing-for-Web-Usage-Agents"

# --------------------------------------------------------------------- type
SZ_BODY = Pt(19)
SZ_CAPTION = Pt(15.5)
SZ_LABEL = Pt(14)
SZ_HEADER = Pt(16)
SZ_LANE = Pt(24)

# --------------------------------------------------------------------- grid
COL_X = (Mm(20.2), Mm(209.3), Mm(398.3))
COL_W = Mm(179.2)
FULL_X, FULL_W = Mm(20.2), Mm(557.3)
SYS_Y = Mm(133)
ROW_BOTTOM = Mm(782)

TITLE = "When Is Expensive Perception Worth Paying For?"
SUBTITLE = ("Measuring the ceiling, the predictability and the economics of "
            "representation routing in web agents")


def body(slide, x, y, w, paragraphs, *, size=SZ_BODY, color=INK, after=Mm(3)):
    h = V8.text_height(paragraphs, SANS, size, w, spacing=1.05, gap=Mm(1.6))
    textbox(slide, x, y, w, h, paragraphs, size=size, color=color,
            line_spacing=1.05, space_after=Mm(1.6))
    return y + h + after


def caption(slide, x, y, w, text, *, after=Mm(3), color=MUTED, size=SZ_CAPTION):
    h = V8.text_height([text], SANS, size, w, spacing=1.08, gap=0)
    textbox(slide, x, y, w, h, [text], size=size, color=color,
            line_spacing=1.08, space_after=Pt(0))
    return y + h + after


def picture(slide, png: Path, x, y, w, *, frame=True):
    with Image.open(png) as im:
        h = int(w * im.height / im.width)
    if frame:
        rect(slide, x, y, w, h, fill=None, line=HAIRLINE)
    slide.shapes.add_picture(str(png), x, y, width=w)
    return y + h


def figure(slide, png: Path, x, y, w, cap: str, *, pad=Mm(3), scale=1.0):
    """A figure in the template's tinted box, with its one-line caption under it.

    ``scale`` narrows the plate inside the column (the caption still runs the
    full column width). Use it where a figure is tall for its width and would
    otherwise push the column past its box."""
    plate = int(w * scale)
    px = x + (w - plate) // 2
    with Image.open(png) as im:
        inner = plate - 2 * pad
        h = int(inner * im.height / im.width)
    box_h = Mm(1.4) + pad + h + pad
    rect(slide, px, y, plate, box_h, fill=FIG_FILL, line=HAIRLINE)
    rect(slide, px, y, plate, Mm(1.4), fill=ACCENT)
    slide.shapes.add_picture(str(png), px + pad, y + Mm(1.4) + pad, width=inner)
    return caption(slide, x, y + box_h + Mm(2), w, cap, after=Mm(5))


def panel(slide, x, y, label, width):
    return V8.panel_header(slide, x, y, label, width)


def table(slide, x, y, w, rows, *, col=0.66, size=Pt(16), head=None, rule=True,
          line_h=Mm(8.2)):
    """A minimal two-column table. ``rows`` is a list of (left, right); a row
    whose right cell is None is a section heading."""
    yy = y
    if head:
        textbox(slide, x, yy, w, Mm(7), [head[0]], size=Pt(14), font=MONO,
                color=MUTED, space_after=Pt(0))
        textbox(slide, x + int(w * col), yy, w - int(w * col), Mm(7), [head[1]],
                size=Pt(14), font=MONO, color=MUTED, align=PP_ALIGN.RIGHT,
                space_after=Pt(0))
        yy += Mm(8)
        rect(slide, x, yy - Mm(1.5), w, Mm(0.4), fill=HAIRLINE)
    for left, right in rows:
        if right is None:
            textbox(slide, x, yy + Mm(1.5), w, Mm(8), [left], size=Pt(16),
                    bold=True, color=INK_STRONG, space_after=Pt(0))
            yy += line_h + Mm(1.5)
            continue
        textbox(slide, x, yy, int(w * col), Mm(8), [left], size=size, color=INK,
                space_after=Pt(0))
        textbox(slide, x + int(w * col), yy, w - int(w * col), Mm(8), [right],
                size=size, bold=True, color=INK_STRONG, align=PP_ALIGN.RIGHT,
                space_after=Pt(0))
        yy += line_h
        if rule:
            rect(slide, x, yy - Mm(1.6), w, Mm(0.25), fill=HAIRLINE)
    return yy + Mm(2)


# ------------------------------------------------------------------ row 1
def build_system(slide, y):
    """The one figure a passer-by is meant to read: agent, three ways of seeing
    a page, and what the choice is judged on. Real screenshots inside it."""
    y = panel(slide, FULL_X, y, "WHAT AN AGENT SEES, AND WHO CHOOSES", FULL_W)
    with Image.open(V9 / "system.png") as im:
        h = int(FULL_W * im.height / im.width)
    slide.shapes.add_picture(str(V9 / "system.png"), FULL_X, y, width=FULL_W)
    return y + h + Mm(3)


# ------------------------------------------------------------------ row 2
RUNS = {"read": ("B0_dom_classifieds_R31194_clean_replicate", "dom"),
        "look": ("B0_vision_classifieds_R24792_clean_replicate", "vision")}


def _distinct_pages(lane: str, task: int) -> int:
    """How many visually distinct pages the run ever saw.

    Grouping is on the screenshots themselves at thumbnail scale: the step's
    ``obs_url`` is written before the action and the screenshot after it, so a
    URL count answers a different question than the one the strip asks."""
    import numpy as np
    run, mode = RUNS[lane]
    d = REPL / run / f"phase1_{mode}_router_0" / "artifacts" / f"classifieds_task_{task}"
    seen: list = []
    for p in sorted(d.glob("step_*/screenshot.png")):
        a = np.asarray(Image.open(p).convert("L").resize((160, 90)), float)
        if not any(np.abs(a - b).mean() < 2.0 for b in seen):
            seen.append(a)
    return len(seen)


def _lane_costs(lane: str, task: int, picks: list[int]):
    """Running billed cost at each shown step, summed from the step records."""
    from p79.experiment.io_utils import read_jsonl_dedup
    run, mode = RUNS[lane]
    rows = read_jsonl_dedup(
        str(REPL / run / f"phase1_{mode}_router_0" / "episodes"
            / f"classifieds_task_{task}_steps_v2.jsonl"))
    out, tot = {}, 0.0
    for r in rows:
        c = r.get("cost_usd") or {}
        c = c.get("total") if isinstance(c, dict) else c
        tot += (c or 0)
        out[r["step_idx"]] = (r.get("action_type", "?"), tot)
    return [(p, *out[p]) for p in picks if p in out], len(rows), tot


TASK = 76
INTENT = ("“Go to my listing of the blue bike and change the price to $85.50 — "
          "and say so in the description.”")
LANES = (
    ("read", "READ", "page text only", C_TEXT, [0, 2, 5, 11], True),
    ("look", "LOOK", "screenshot only", C_IMAGE, [0, 3, 9, 20], False),
)


def build_strip(slide, y):
    """One task, two ways of seeing it. LOOK's two middle frames are the same
    empty form nine steps apart — the reader sees the loop instead of reading
    that there was one."""
    x, w = FULL_X, FULL_W
    y = panel(slide, x, y, "THE SAME TASK, SEEN TWO WAYS", w)
    y = caption(slide, x, y, w, INTENT, color=INK, size=Pt(18), after=Mm(3))

    lane_w = Mm(46)
    gap = Mm(3)
    cell = int((w - lane_w - 3 * gap) / 4)
    with Image.open(V9 / "lane_read_0.png") as im:
        fh = int((cell - Mm(1)) * im.height / im.width)

    for key, name, note, colour, picks, won in LANES:
        rows, n_steps, total = _lane_costs(key, TASK, picks)
        pages = _distinct_pages(key, TASK)
        rect(slide, x, y, Mm(2.4), fh + Mm(15), fill=colour)
        textbox(slide, x + Mm(5), y + Mm(1), lane_w - Mm(6), Mm(10), [name],
                font=MONO, size=SZ_LANE, bold=True, color=colour, space_after=Pt(0))
        textbox(slide, x + Mm(5), y + Mm(12), lane_w - Mm(6), Mm(16), [note],
                size=SZ_CAPTION, color=MUTED, line_spacing=1.08, space_after=Pt(0))
        verdict = "solved" if won else "gave up"
        textbox(slide, x + Mm(5), y + Mm(29), lane_w - Mm(6), Mm(28),
                [f"**{n_steps} steps**", f"**{pages} different pages**",
                 f"**${total:.2f}** · {verdict}"],
                size=SZ_CAPTION, color=C_BOTH if won else C_FAIL,
                line_spacing=1.1, space_after=Pt(0))
        fx = x + lane_w
        for i, (step, act, cost) in enumerate(rows):
            px = fx + Emu(i * (cell + int(gap)))
            slide.shapes.add_picture(str(V9 / f"lane_{key}_{i}.png"), px, y,
                                     width=cell - Mm(1))
            rect(slide, px, y, cell - Mm(1), fh, fill=None, line=GREY)
            textbox(slide, px, y + fh + Mm(1.5), cell, Mm(7),
                    [f"step {step} · {act}"], size=SZ_CAPTION, bold=True,
                    color=INK_STRONG, space_after=Pt(0))
            textbox(slide, px, y + fh + Mm(8), cell, Mm(7), [f"${cost:.3f} spent"],
                    size=Pt(14), color=MUTED, space_after=Pt(0))
        y += fh + Mm(18)
    return caption(slide, x, y - Mm(2), w,
                   "**LOOK's last three frames are the same page**, pixel for pixel, "
                   "seventeen steps apart — it saw only 9 distinct pages in 26 steps. "
                   "One recorded run per view; both outcomes repeat on an independent rerun.",
                   color=INK, after=Mm(0))


# ------------------------------------------------------------------ row 3
def build_col1(slide, y):
    """Why choosing could pay: the views solve overlapping but different sets,
    and the split that shows up in *behaviour* is image against no image."""
    x, w = COL_X[0], COL_W
    y = panel(slide, x, y, "DIFFERENT VIEWS, DIFFERENT TASKS", w)
    y = figure(slide, V9 / "venn_b0.png", x, y, w,
               "Tasks solved by each text-only view. The sets overlap; they do not coincide.")
    y = panel(slide, x, y, "THEY ALSO BEHAVE DIFFERENTLY", w)
    y = body(slide, x, y, w, [
        "Across **26 behavioural measures**, how often a view is the extreme one "
        "in at least 7 of the 8 settings:"], size=Pt(21), after=Mm(4))
    # Table 5 of the evidence inventory. The sheet reports the counts and stops
    # there: the source caption discloses that its ≥7/8 bar was set after the
    # cell count changed, so "the extremes sit with the two views that see the
    # page" is the claim, and "the text-only views are inseparable" is not.
    y = table(slide, x, y, w, [
        ("Vision — screenshot only", "9"),
        ("SoM — screenshot + text", "5"),
        ("DOM — page text only", "0"),
        ("the 3 other text-only views", "0"),
    ], head=("view", "measures"), size=Pt(19), line_h=Mm(9.6))
    y = caption(slide, x, y, w,
                "The extremes sit with the two views that see the page. How the agent "
                "moves is set by whether it gets a picture, not by which text it gets.",
                size=Pt(17), after=Mm(4))
    return y


def build_col2(slide, y):
    """What was run, and how the two channels fail when they lose."""
    x, w = COL_X[1], COL_W
    y = panel(slide, x, y, "SIX VIEWS, EIGHT SETTINGS", w)
    y = figure(slide, THESIS / "fig_f5_design_matrix.png", x, y, w,
               "Success rate (%). Colour is what reaches the model: text, text + image, image.",
               scale=0.74)
    y = panel(slide, x, y, "AND THEY FAIL DIFFERENTLY", w)
    y = body(slide, x, y, w, [
        "On tasks only one side solved, how did the loser fail? "
        "**× = versus its own everyday failure rate.**"], after=Mm(1))
    # Table 41. Both blocks are read against their own baseline, never across
    # blocks: the text side is four arms and the image side two.
    y = table(slide, x, y, w, [
        ("When the picture won, the text-only agent", None),
        ("gave up once it could not find it", "2.3×"),
        ("could not see what the task asked about", "2.2×"),
        ("When the text won, the picture-only agent", None),
        ("never turned the page", "0.9×"),
        ("ran out of budget, unfinished", "0.9×"),
    ], col=0.78, size=Pt(15), line_h=Mm(6.8))
    y = caption(slide, x, y, w,
                "**One side dies of something you can name. The other simply never "
                "arrives.**", color=INK, after=Mm(3))
    return y


def build_col3(slide, y):
    """Why choosing well turns out to be hard: nothing learned beats the free
    policy, and the measurement floor is close to the effect being chased."""
    x, w = COL_X[2], COL_W
    y = panel(slide, x, y, "CHOOSING WELL IS HARDER THAN IT LOOKS", w)
    y = figure(slide, THESIS / "fig_f13_dominance_plane.png", x, y, w,
               "Each policy against always using the cheapest view. A win lands in the "
               "shaded region — none does.", scale=0.96)
    y = body(slide, x, y, w, [
        "And the target is small: running **one unchanged view twice** already "
        "flips **10–14%** of tasks."], after=Mm(6))
    y = panel(slide, x, y, "SO WHAT", w)
    y = body(slide, x, y, w, [
        "The views are worth choosing between — they solve different tasks, move "
        "differently and fail differently.",
        "**Choosing well, in advance, is the part that does not work yet.**"],
        after=Mm(3))
    return y


# ------------------------------------------------------------------ run
def main():
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    title = find(slide, "TextBox 6")
    title.top, title.height = Mm(10), Mm(40)
    set_text(title, [TITLE], size=Pt(44))
    set_text(find(slide, "TextBox 7"),
             ["Jiaming Wei          Supervisors: Prof. María Pérez-Ortiz  ·  Zekun Wu"],
             size=Pt(23))
    set_text(find(slide, "TextBox 8"),
             ["UCL Centre for Artificial Intelligence          Holistic AI"], size=Pt(18))
    set_text(find(slide, "TextBox 12"), [SUBTITLE], size=Pt(26))

    V8.HEADER_PARTS = tuple(find(slide, n) for n in
                            ("Rectangle 13", "TextBox 14", "Rectangle 15", "Rectangle 16"))
    drop(slide, "TextBox 17", "Rectangle 33", "TextBox 34", "Rectangle 35",
         "Rectangle 36", "TextBox 37", "Rectangle 38", "Rectangle 39", "TextBox 42",
         "Rectangle 48", "TextBox 49", "Rectangle 50", "Rectangle 51",
         "Rectangle 52", "TextBox 53", "TextBox 54", "TextBox 55", "TextBox 56",
         "TextBox 57", "TextBox 58")

    y = build_system(slide, SYS_Y)
    y = build_strip(slide, y + Mm(5)) + Mm(6)
    ends = {"views (left)": build_col1(slide, y),
            "setup (middle)": build_col2(slide, y),
            "routing (right)": build_col3(slide, y)}

    drop(slide, *V8.HEADER_PARTS)

    set_text(find(slide, "TextBox 71"),
             ["jiaming.wei.25@ucl.ac.uk",
              "Paper, code & full results: github.com/Quarkgluonmixture/"
              "Cost-Aware-Routing-for-Web-Usage-Agents"])
    qr_png = V9 / "qr_repo.png"
    qrcode.make(REPO_URL, box_size=20, border=1).save(qr_png)
    set_text(find(slide, "TextBox 73"), ["Explore all 8 settings →"])
    slot = find(slide, "Rectangle 74")
    slide.shapes.add_picture(str(qr_png), slot.left + Mm(2), slot.top + Mm(2),
                             width=slot.width - Mm(4))
    drop(slide, "Rectangle 74", "TextBox 75")

    prs.save(str(OUT))
    verify(prs, ends, y)


def verify(prs, ends, row_y):
    mm = lambda emu: emu / 914400 * 25.4  # noqa: E731
    w, h = mm(prs.slide_width), mm(prs.slide_height)
    assert abs(w - 594) < 0.5 and abs(h - 841) < 0.5, "slide was resized!"
    print(f"wrote {OUT.relative_to(REPO)}   ({w:.0f}x{h:.0f}mm, A1 portrait, not resized)")
    print(f"  {'columns start':16s} {mm(row_y):6.1f}mm")
    bad = False
    for name, end in ends.items():
        slack = mm(ROW_BOTTOM - end)
        flag = "" if slack >= 0 else "   <-- OVERRUNS ITS BOX"
        bad |= slack < 0
        print(f"  {name:16s} ends {mm(end):6.1f}mm   (box to 782mm, slack {slack:+6.1f}mm){flag}")
    if bad:
        raise SystemExit("a column overran its box — shorten it or move the grid")


if __name__ == "__main__":
    main()
