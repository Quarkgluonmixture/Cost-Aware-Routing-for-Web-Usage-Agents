#!/usr/bin/env python3
"""Build the Holistic AI x UCL CDI showcase poster from the supplied A1 template.

v8 — v6's look, v7's words (2026-09-03)
----------------------------------------
v7 moved the loop diagram to the foot and put text first; the author: "the
previous one looked better". It did — the top half had become prose. v8 puts
the loop back on top at v6 size, the number strip and the three comparison
definitions under it, and keeps every v7 wording fix.
The poster now stands next to a laptop that replays the same task through three
ways of seeing the page. That split decides what goes on silk:

    the DEMO shows the phenomenon   — one task, three eyes, three behaviours,
                                      three bills, step by step
    the POSTER shows the system and — where the decision sits in the agent loop,
    the measurement                   what it was worth in hindsight, whether it
                                      could be learned, and why not

Reviewer brief (Holistic AI, 2026-09-02) still holds: one main system diagram,
compact, one type scale, the template's own skeleton. So:

    TITLE        Look, read, or both? / Web agents can't yet learn how to see a page
    STANDFIRST   one sentence, both baselines named, two lines at the template's 28pt
    FIG 1        the system diagram, drawn in native shapes at full width:
                 task + page -> WHO DECIDES HOW TO SEE? -> LOOK / READ / BOTH ->
                 agent step -> outcome + bill. The decision box is the experiment.
    COLUMN 1     ON THE SCREEN BESIDE YOU: the three demo tasks, one frame each,
                 with every mode's real outcome / steps / bill
    COLUMNS 2-3  RESULTS: Fig 2 (thesis F13) + three verdicts; WHY IT CANNOT BE
                 LEARNED; TAKEAWAY

Type scale is the template's own (read from its placeholders):
    title 41 Georgia · standfirst 28 Georgia · section header 14 Consolas bold
    body 17.65 Arial · caption 12.7 Arial grey · marks 32.5 Consolas bold

Every number on the demo strip is parsed from the episode summaries by
``poster_figures.py`` (``figures/demo_strip.json``) — nothing on that strip is
typed by hand, and a task whose outcome flips on rerun fails the build there.

Hard constraint from the organisers: **do not resize the slide.** ``verify``
re-asserts it and fails the build if any panel overruns its box.

Usage::

    .venv/bin/python3 deliverables/showcase/poster_figures.py   # figures + strip data
    .venv/bin/python3 deliverables/showcase/build_poster.py
"""

from __future__ import annotations

import copy
import json
import re
from pathlib import Path

import qrcode
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Mm, Pt

REPO = Path(__file__).resolve().parents[2]
HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "Showcase-Poster-Template-A1.pptx"
OUT = HERE / "poster_jiaming_wei.pptx"
FIGDIR = HERE / "figures"

REPO_URL = "https://github.com/Quarkgluonmixture/Cost-Aware-Routing-for-Web-Usage-Agents"

# ---------------------------------------------------------------- palette / type
INK = RGBColor(0x12, 0x16, 0x2E)
INK_STRONG = RGBColor(0x14, 0x1E, 0x41)
MUTED = RGBColor(0x5D, 0x67, 0x87)
ACCENT = RGBColor(0x50, 0x49, 0xF9)
HAIRLINE = RGBColor(0xDD, 0xE3, 0xF2)
FIG_FILL = RGBColor(0xF7, 0xF9, 0xFD)
GREY = RGBColor(0xC9, 0xCF, 0xE0)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
# Thesis figure palette: orange = text only, green = text + image, blue = image
# only. The demo uses the same three colours for READ / BOTH / LOOK.
C_TEXT = RGBColor(0xE8, 0x72, 0x0C)
C_BOTH = RGBColor(0x14, 0x85, 0x5F)
C_IMAGE = RGBColor(0x1F, 0x5F, 0xD6)
C_FAIL = RGBColor(0xC2, 0x35, 0x2B)

SERIF, SANS, MONO = "Georgia", "Arial", "Consolas"

FONT_FILE = {
    (SANS, False): "/usr/share/fonts/truetype/croscore/Arimo-Regular.ttf",
    (SANS, True): "/usr/share/fonts/truetype/croscore/Arimo-Bold.ttf",
    (SERIF, False): "/usr/share/fonts/truetype/noto/NotoSerif-Regular.ttf",
    (SERIF, True): "/usr/share/fonts/truetype/noto/NotoSerif-Bold.ttf",
}
# Calibrated to the renderer (LibreOffice), not the face file: see §498.2.
LINE_HEIGHT = {SANS: 1.20, SERIF: 1.20}

SZ_TITLE = Pt(40.96)
SZ_STANDFIRST = Pt(28)
SZ_BODY = Pt(17.65)
SZ_CAPTION = Pt(12.71)
SZ_MARK = Pt(32.48)

BODY_SPACING = 1.2
CAPTION_SPACING = 1.32
PARA_GAP = Mm(3.5)

# ------------------------------------------------------------------------- grid
COL_X = (Mm(20.2), Mm(209.3), Mm(398.3))
COL_W = Mm(179.2)
SPAN_23_X, SPAN_23_W = Mm(209.3), Mm(368.2)
FULL_X, FULL_W = Mm(20.2), Mm(557.3)

LOOP_Y = Mm(140)         # the agent loop, full width, on top (the hero)
ROW_BOTTOM = Mm(782)

MODES = (("look", "LOOK", C_IMAGE), ("read", "READ", C_TEXT), ("both", "BOTH", C_BOTH))


# ------------------------------------------------------------------- measurement
_MARK_RE = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*")
_SCALE = 8


def segments(text: str):
    out, pos = [], 0
    for m in _MARK_RE.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], False, False))
        if m.group(1) is not None:
            out.append((m.group(1), True, False))
        else:
            out.append((m.group(2), False, True))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False, False))
    out = out or [("", False, False)]
    leftover = "".join(chunk for chunk, _, _ in out)
    assert "*" not in leftover, f"unpaired '*' would print literally: {text!r}"
    return out


_font_cache: dict = {}


def _font(family: str, bold: bool, size_pt: float):
    key = (family, bold, round(size_pt, 2))
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(
            FONT_FILE[(family, bold)], int(round(size_pt * _SCALE)))
    return _font_cache[key]


def line_count(text: str, family: str, size_pt: float, width_emu: int) -> int:
    width_px = width_emu / 12700 * _SCALE
    words = [(w, b) for seg, b, _ in segments(text) for w in seg.split(" ") if w]
    lines, cur_px, started = 1, 0.0, False
    for word, bold in words:
        f = _font(family, bold, size_pt)
        w_px = f.getlength(word if not started else " " + word)
        if started and cur_px + w_px > width_px:
            lines += 1
            cur_px = f.getlength(word)
        else:
            cur_px += w_px
            started = True
    return lines


def text_height(paragraphs, family, size, width_emu, *, spacing, gap) -> int:
    n = sum(line_count(p, family, size.pt, width_emu) for p in paragraphs)
    return int(n * size.pt * LINE_HEIGHT[family] * spacing * 12700) + gap * (len(paragraphs) - 1)


# ------------------------------------------------------------------- pptx helpers
def find(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"shape {name!r} not in template")


def drop(slide, *shapes):
    for shape in shapes:
        el = shape._element if hasattr(shape, "_element") else find(slide, shape)._element
        el.getparent().remove(el)


def clone(slide, shape, left, top):
    el = copy.deepcopy(shape._element)
    shape._element.getparent().append(el)
    new = slide.shapes[-1]
    new.left, new.top = left, top
    return new


def _run(paragraph, text, font, size, color, bold, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _emit(paragraph, text, *, font, size, color, bold=False):
    for chunk, is_bold, is_italic in segments(text):
        _run(paragraph, chunk, font, size, color, bold or is_bold, is_italic)


def textbox(slide, left, top, width, height, paragraphs, *, font=SANS, size=SZ_BODY,
            color=INK, line_spacing=BODY_SPACING, space_after=PARA_GAP,
            align=PP_ALIGN.LEFT, bold=False, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    if anchor is not None:
        frame.vertical_anchor = anchor
    for i, text in enumerate(paragraphs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.line_spacing = line_spacing
        para.space_after = space_after
        para.alignment = align
        _emit(para, text, font=font, size=size, color=color, bold=bold)
    return box


def body(slide, x, y, w, paragraphs, *, after=Mm(4), size=SZ_BODY, color=INK):
    h = text_height(paragraphs, SANS, size, w, spacing=BODY_SPACING, gap=PARA_GAP)
    textbox(slide, x, y, w, h, paragraphs, size=size, color=color)
    return y + h + after


def caption(slide, x, y, w, paragraphs, *, after=Mm(4), color=MUTED):
    h = text_height(paragraphs, SANS, SZ_CAPTION, w, spacing=CAPTION_SPACING, gap=Mm(1.5))
    textbox(slide, x, y, w, h, paragraphs, size=SZ_CAPTION, color=color,
            line_spacing=CAPTION_SPACING, space_after=Mm(1.5))
    return y + h + after


def set_text(shape, paragraphs, *, size=None, align=None):
    frame = shape.text_frame
    first = frame.paragraphs[0]
    proto = first.runs[0]
    font, color, bold = proto.font.name, proto.font.color.rgb, proto.font.bold
    size = size or proto.font.size
    spacing = first.line_spacing
    frame.clear()
    for i, text in enumerate(paragraphs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.line_spacing = spacing
        if align is not None:
            para.alignment = align
        _emit(para, text, font=font, size=size, color=color, bold=bool(bold))


def rect(slide, x, y, w, h, fill=None, line=None, *, lw=Pt(0.7), radius=None, dash=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = lw
        if dash:
            shape.line.dash_style = MSO_LINE.DASH
    shape.shadow.inherit = False
    if radius is not None:
        shape.adjustments[0] = radius
    shape.text_frame.word_wrap = True
    return shape


def arrow_right(slide, x, y, w, h, fill=GREY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


HEADER_PARTS: tuple = ()


def panel_header(slide, x, y, label, width):
    block, title, rule_a, rule_b = HEADER_PARTS
    base = block.top
    clone(slide, block, x, y)
    t = clone(slide, title, x + (title.left - block.left), y - Mm(1.2))
    ra = clone(slide, rule_a, x, y + (rule_a.top - base))
    rb = clone(slide, rule_b, x + (rule_b.left - block.left), y + (rule_b.top - base))
    rb.width = width - ra.width - (rb.left - x - ra.width)
    set_text(t, [label])
    return y + Mm(15)


def fig_box(slide, x, y, w, png: Path, cap: str):
    inset = Mm(3)
    with Image.open(png) as im:
        pic_w = w - 2 * inset
        pic_h = int(pic_w * im.height / im.width)
    box_h = Mm(1.4) + inset + pic_h + inset
    rect(slide, x, y, w, box_h, fill=FIG_FILL, line=HAIRLINE)
    rect(slide, x, y, w, Mm(1.4), fill=ACCENT)
    slide.shapes.add_picture(str(png), x + inset, y + Mm(1.4) + inset, width=pic_w)
    return caption(slide, x, y + box_h + Mm(2.5), w, [cap], after=Mm(3))


# ----------------------------------------------------------------- primitives
def card(slide, x, y, w, h, *, bar=ACCENT, dash=False, line=HAIRLINE):
    rect(slide, x, y, w, h, fill=FIG_FILL, line=line, dash=dash)
    if bar is not None:
        rect(slide, x, y, w, Mm(1.4), fill=bar)


def label(slide, x, y, w, text, color=MUTED):
    textbox(slide, x, y, w, Mm(6), [text], font=MONO, size=Pt(10.59), color=color,
            space_after=Pt(0))
    return y + Mm(6.5)


def _pic(slide, png: Path, x, y, w, *, frame=True):
    with Image.open(png) as im:
        h = int(w * im.height / im.width)
    if frame:
        rect(slide, x - Mm(0.4), y - Mm(0.4), w + Mm(0.8), h + Mm(0.8), fill=None, line=GREY)
    slide.shapes.add_picture(str(png), x, y, width=w)
    return h


def _read_lines() -> list[str]:
    out = []
    for ln in (FIGDIR / "eye_read.txt").read_text(encoding="utf-8").splitlines():
        ln = ln.strip().split(" url:")[0].split(" focused:")[0]
        out.append(ln[:36])
    return out[:6]


def metric_strip(slide, x, y, w, tiles):
    h = Mm(32.5)
    rect(slide, x, y, w, h, fill=RGBColor(0xEC, 0xEF, 0xFA))
    pad = Mm(7)
    tile_w = int((w - 2 * pad) / len(tiles))
    for i, (number, lbl) in enumerate(tiles):
        assert len(lbl) * 10.59 * 0.6 * 12700 < tile_w - Mm(4), f"metric label overprints: {lbl!r}"
        tx = x + pad + Emu(i * tile_w)
        textbox(slide, tx, y + Mm(5), Emu(tile_w), Mm(17), [number], font=MONO,
                size=SZ_MARK, color=INK_STRONG, bold=True, line_spacing=1.18, space_after=Pt(0))
        textbox(slide, tx, y + Mm(20.5), Emu(tile_w), Mm(9), [lbl], font=MONO,
                size=Pt(10.59), color=MUTED, line_spacing=1.18, space_after=Pt(0))
    return y + h + Mm(4)


# ------------------------------------------------------------ row A: the result
def build_result_strip(slide, y):
    """The three numbers, and the three comparisons they are made against,
    defined in one line each. The hindsight line carries its own health
    warning (the rerun band)."""
    x, w = FULL_X, FULL_W
    y = panel_header(slide, x, y, "RESULTS ACROSS 8,934 TASK ATTEMPTS", w)
    y = metric_strip(slide, x, y, w, [
        ("+16.35 in 100", "PERFECT HINDSIGHT, BEST OF 8 · VS ONE FIXED VIEW"),
        ("0 of 8", "LEARNED CHOICES THAT BEAT ALWAYS-CHEAPEST"),
        ("1 of 8", "HINDSIGHT CHOICES THAT BEAT ALWAYS-CHEAPEST"),
    ])
    keys = [
        "**ONE FIXED VIEW** — the single view that solves most tasks in a setting, "
        "used for every task.",
        "**ALWAYS-CHEAPEST** — the single view that costs least on average in a "
        "setting, used for every task.",
        "**PERFECT HINDSIGHT** — for each task, the view that solved it, picked after "
        "the fact. An optimistic bound: rerunning one unchanged view flips 10–14% of "
        "outcomes and by itself gains 2.0–7.6 tasks in 100.",
    ]
    kw = int((w - 2 * Mm(6)) / 3)
    ends = []
    for i, text in enumerate(keys):
        ends.append(caption(slide, x + Emu(i * (kw + int(Mm(6)))), y, Emu(kw), [text],
                            color=INK, after=Mm(0)))
    return max(ends)


# ------------------------------------------------------ row B left: the why
def build_why(slide, y):
    x, w = COL_X[0], COL_W
    # the catch, as a callout in the template's tinted-box idiom
    line = ("**The agents that would gain most from choosing a view solve the fewest "
            "tasks — so they produce the fewest examples to learn the choice from.**")
    h = text_height([line], SANS, SZ_BODY, w - Mm(12), spacing=BODY_SPACING, gap=0)
    box_h = Mm(6) + Mm(6.5) + h + Mm(6)
    rect(slide, x, y, w, box_h, fill=RGBColor(0xEC, 0xEF, 0xFA))
    rect(slide, x, y, Mm(2.4), box_h, fill=ACCENT)
    yy = label(slide, x + Mm(6), y + Mm(5), w - Mm(12), "THE CATCH", color=ACCENT)
    textbox(slide, x + Mm(6), yy, w - Mm(12), h, [line])
    y += box_h + Mm(8)

    y = panel_header(slide, x, y, "WHY LEARNING THE CHOICE FAILS HERE", w)
    y = body(slide, x, y, w, [
        "A training example for “which view” exists only when the agent solves a "
        "task. Here the best single view solves just **2–36%** of tasks, leaving "
        "typically **15–97** usable examples per setting — enough to train a "
        "classifier in only 2 of the 6 VisualWebArena settings."], after=Mm(4))
    y = fig_box(slide, x, y, w, FIGDIR / "poster_label_supply.png",
                "**Fig 3.** Usable “which view” examples against the best single view's "
                "success rate, one point per VisualWebArena setting. Examples exist only "
                "where tasks get solved.")
    y = body(slide, x, y, w, [
        "Shrinking the training data on purpose points to scarcity as the main "
        "bottleneck, and prices it: the failing settings would need at least "
        "**2.1–4.2×** more tasks than the benchmarks contain."], after=Mm(7))

    y = panel_header(slide, x, y, "TAKEAWAY", w)
    y = body(slide, x, y, w, [
        "**Learning how to see is not only a modelling problem: whether it can be "
        "learned depends on how good the agent producing the examples already is.** "
        "So, in this order: improve the agent, then collect reliable examples, then "
        "learn when to look."], after=Mm(1.5))
    return caption(slide, x, y, w, [
        "Measured in the 2–36% success regime we observed. This need not hold for "
        "stronger agents."], after=Mm(0))


# ------------------------------------------------- row B right: the evidence
def build_evidence(slide, y):
    x, w = SPAN_23_X, SPAN_23_W
    y = fig_box(
        slide, x, y, w, FIGDIR / "poster_dominance_plane.png",
        "**Fig 2.** Every way of choosing a view, in every setting, compared with one "
        "fixed rule: **always use the cheapest view** (★). A win lands in the shaded "
        "region — cheaper *and* no worse. Always-cheapest is cheapest on average, not "
        "on every task, which is why a few points sit left of it. Learned choices are "
        "scored only on tasks they never saw.")
    y = body(slide, x, y, w, [
        "**Perfect hindsight would pay.** Picking the winning view for each task after "
        "the fact would solve **3.45 to 16.35 more tasks in every 100** than the best "
        "single view, and spend 1.6–35.3% less — in all 8 settings.",
        "**Nothing we trained could do it.** In **0 of 8** settings did a learned "
        "choice beat always using the cheapest view on both success and cost — and "
        "even perfect hindsight manages that in only **1 of 8**.",
        "**What survives is a bound, not a method.** Sending the tasks nobody solves to "
        "the cheapest view saves 9.5–30.6% at the same success in 8 of 8 — but that "
        "too needs hindsight, and plain always-cheapest usually saves more."],
        after=Mm(5))

    # the bridge to the laptop: results only, no frames — the screen has those
    y = panel_header(slide, x, y, "WATCH THE LAPTOP BESIDE THIS POSTER", w)
    y = body(slide, x, y, w, [
        "Same task, three views, different behaviour and different bills — step by "
        "step. Three illustrative tasks, chosen so each view wins once; **not** how "
        "often each wins."], after=Mm(3))
    strip = json.loads((FIGDIR / "demo_strip.json").read_text())
    task_w = Mm(200)
    cell_w = int((w - task_w) / 3)
    for i, (key, name, colour) in enumerate(MODES):
        textbox(slide, x + task_w + Emu(i * cell_w), y, Emu(cell_w), Mm(6), [name],
                font=MONO, size=Pt(10.59), bold=True, color=colour, space_after=Pt(0))
    y += Mm(7)
    rect(slide, x, y, w, Mm(0.3), fill=HAIRLINE)
    y += Mm(2.5)
    for task, d in strip.items():
        intent = f"“{d['intent']}”"
        ih = text_height([intent], SANS, SZ_CAPTION, task_w - Mm(6), spacing=CAPTION_SPACING, gap=0)
        textbox(slide, x, y, task_w - Mm(6), ih, [intent], size=SZ_CAPTION, color=INK,
                line_spacing=CAPTION_SPACING, space_after=Pt(0))
        for i, (key, name, colour) in enumerate(MODES):
            m = d["modes"][key]
            mark, mc = ("✓", C_BOTH) if m["success"] else ("✗", C_FAIL)
            box = slide.shapes.add_textbox(x + task_w + Emu(i * cell_w), y, Emu(cell_w), Mm(8))
            fr = box.text_frame; fr.word_wrap = True
            fr.margin_left = fr.margin_right = fr.margin_top = fr.margin_bottom = 0
            para = fr.paragraphs[0]
            _run(para, mark + " ", MONO, Pt(15), mc, True)
            _run(para, f"{m['steps']} steps · ", SANS, SZ_CAPTION, INK, False)
            _run(para, f"${m['cost_usd']:.3f}", SANS, SZ_CAPTION, INK, True)
        y += max(ih, int(Mm(8))) + Mm(3)
        rect(slide, x, y - Mm(1), w, Mm(0.3), fill=HAIRLINE)
    return caption(slide, x, y + Mm(1), w, [
        "One recorded attempt per view, B0 · classifieds. Every ✓ / ✗ came out the "
        "same on an independent rerun; steps and bill differ run to run."], after=Mm(0))


# --------------------------------------------------- row C: the agent loop
def build_loop(slide):
    """Fig 1, at the foot: how the measurement was made, as a loop that loops."""
    y = panel_header(slide, FULL_X, LOOP_Y, "THE AGENT LOOP, AND WHERE THE DECISION SITS", FULL_W)
    top, H = y, Mm(140)
    x0 = FULL_X
    pad = Mm(5)
    cols = {"page": (0, 84), "decide": (98, 108), "eyes": (220, 150), "agent": (384, 88),
            "stop": (486, 71.3)}
    for ax in (86, 208, 372, 474):
        arrow_right(slide, x0 + Mm(ax), top + H / 2 - Mm(7), Mm(10), Mm(14))

    cx, cw = x0 + Mm(cols["page"][0]), Mm(cols["page"][1])
    card(slide, cx, top, cw, H)
    yy = label(slide, cx + pad, top + Mm(4), cw - 2 * pad, "THE TASK + THE LIVE PAGE")
    yy = body(slide, cx + pad, yy, cw - 2 * pad,
              ["“Show me the cheapest bike with red handlebars between $900–950.”"],
              after=Mm(2.5))
    yy = caption(slide, cx + pad, yy, cw - 2 * pad,
                 ["Part of the intent is in the pictures, part in the text."], after=Mm(3))
    _pic(slide, FIGDIR / "eye_look.png", cx + pad, yy, cw - 2 * pad)

    cx, cw = x0 + Mm(cols["decide"][0]), Mm(cols["decide"][1])
    card(slide, cx, top, cw, H, bar=None, dash=True, line=ACCENT)
    yy = label(slide, cx + pad, top + Mm(4), cw - 2 * pad, "WHO DECIDES HOW TO SEE IT?", color=ACCENT)
    pills = [("One fixed view", "the same view for every task"),
             ("Perfect hindsight", "knowing afterwards which view solved it"),
             ("A learned choice", "made before the task runs, from what the page looks like")]
    ph = Mm(34)
    for i, (name, note) in enumerate(pills):
        py = yy + Emu(i * int(ph + Mm(3)))
        rect(slide, cx + pad, py, cw - 2 * pad, ph, fill=PAPER, line=HAIRLINE, radius=0.12)
        textbox(slide, cx + pad + Mm(4), py + Mm(3.5), cw - 2 * pad - Mm(8), Mm(8),
                [name], bold=True, color=INK_STRONG, space_after=Pt(0))
        textbox(slide, cx + pad + Mm(4), py + Mm(13), cw - 2 * pad - Mm(8), Mm(16),
                [note], size=SZ_CAPTION, color=MUTED, line_spacing=CAPTION_SPACING,
                space_after=Pt(0))
    textbox(slide, cx + pad, top + H - Mm(10), cw - 2 * pad, Mm(6),
            ["this box is what we measure"], font=MONO, size=Pt(10.59), color=ACCENT,
            space_after=Pt(0))

    cx, cw = x0 + Mm(cols["eyes"][0]), Mm(cols["eyes"][1])
    eyes = [("LOOK", C_IMAGE, "the screenshot only", "3,123 tokens on this page", "eye_look.png"),
            ("READ", C_TEXT, "the page as text: its elements and labels",
             "3,314 tokens · +3 text-only variants", None),
            ("BOTH", C_BOTH, "the screenshot with numbered boxes, plus the text",
             "4,335 tokens", "eye_both.png")]
    eh, thumb_w = Mm(44), Mm(66)
    for i, (name, colour, what, price, png) in enumerate(eyes):
        ey = top + Emu(i * int(eh + Mm(4)))
        card(slide, cx, ey, cw, eh, bar=colour)
        tx, tw = cx + pad, cw - 2 * pad - thumb_w - Mm(4)
        textbox(slide, tx, ey + Mm(4.5), Mm(40), Mm(8), [name], font=MONO,
                size=SZ_BODY, bold=True, color=colour, space_after=Pt(0))
        textbox(slide, tx, ey + Mm(14), tw, Mm(14), [what], size=SZ_CAPTION, bold=True,
                color=INK_STRONG, line_spacing=CAPTION_SPACING, space_after=Pt(0))
        textbox(slide, tx, ey + Mm(29), tw, Mm(12), [price], size=SZ_CAPTION, color=MUTED,
                line_spacing=CAPTION_SPACING, space_after=Pt(0))
        px, py = cx + cw - pad - thumb_w, ey + Mm(4)
        if png:
            _pic(slide, FIGDIR / png, px, py, thumb_w)
        else:
            rect(slide, px, py, thumb_w, Mm(36.9), fill=PAPER, line=GREY)
            textbox(slide, px + Mm(2), py + Mm(2), thumb_w - Mm(4), Mm(33), _read_lines(),
                    font=MONO, size=Pt(7.5), color=INK, line_spacing=1.15, space_after=Pt(0))

    cx, cw = x0 + Mm(cols["agent"][0]), Mm(cols["agent"][1])
    card(slide, cx, top, cw, H)
    yy = label(slide, cx + pad, top + Mm(4), cw - 2 * pad, "THE AGENT, ONE STEP AT A TIME")
    yy = body(slide, cx + pad, yy, cw - 2 * pad, ["**think → act**"], after=Mm(2))
    yy = caption(slide, cx + pad, yy, cw - 2 * pad,
                 ["click · type · scroll · go back · finish"], after=Mm(4))
    caption(slide, cx + pad, yy, cw - 2 * pad,
            ["Only the page view changes; model, prompt, step budget and cost "
             "accounting stay fixed."], color=INK)

    cx, cw = x0 + Mm(cols["stop"][0]), Mm(cols["stop"][1])
    card(slide, cx, top, cw, H)
    yy = label(slide, cx + pad, top + Mm(4), cw - 2 * pad, "WHEN IT STOPS")
    yy = body(slide, cx + pad, yy, cw - 2 * pad, ["**✓ solved / ✗ not**", "**$ for the attempt**"],
              after=Mm(2))
    caption(slide, cx + pad, yy, cw - 2 * pad,
            ["8 website × model combinations · two public benchmarks · 6 views · "
             "8,934 attempts"])

    agent_cx = x0 + Mm(cols["agent"][0] + cols["agent"][1] / 2)
    page_cx = x0 + Mm(cols["page"][0] + cols["page"][1] / 2)
    ly = top + H + Mm(12)
    rect(slide, agent_cx - Mm(0.6), top + H, Mm(1.2), ly - (top + H), fill=GREY)
    rect(slide, page_cx, ly - Mm(0.6), agent_cx - page_cx, Mm(1.2), fill=GREY)
    rect(slide, page_cx - Mm(0.6), top + H + Mm(4), Mm(1.2), ly - (top + H) - Mm(4), fill=GREY)
    head = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, page_cx - Mm(3), top + H,
                                  Mm(6), Mm(4.5))
    head.fill.solid(); head.fill.fore_color.rgb = GREY; head.line.fill.background()
    head.shadow.inherit = False
    textbox(slide, page_cx + Mm(10), ly + Mm(2), agent_cx - page_cx - Mm(20), Mm(7),
            ["the action changes the page → next step · up to 30 steps per task · "
             "the bill grows with every step"], size=SZ_CAPTION, color=MUTED,
            align=PP_ALIGN.CENTER, space_after=Pt(0))
    return caption(slide, FULL_X, ly + Mm(10), FULL_W, [
        "**Fig 1.** Everything is held fixed except how the page is shown. The dashed box "
        "is what this work measures — one fixed view, perfect hindsight, or a choice a "
        "model had to learn — each judged on **both** success and cost. The laptop "
        "beside this poster replays the loop on three tasks."], after=Mm(0))


# --------------------------------------------------------------------------- run
TITLE = ["Look, read, or both?", "Today's web agents can't yet learn how to see a page"]
STANDFIRST = (
    "Perfect hindsight would solve up to 16 more tasks in 100 than the best single "
    "view. None of 8 learned choices beat always using the cheapest view on both "
    "success and cost."
)


def main():
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    title = find(slide, "TextBox 6")
    title.top, title.height = Mm(8), Mm(46)
    for line in TITLE:
        n = line_count(line, SERIF, SZ_TITLE.pt, int(title.width * 0.94))
        assert n == 1, f"title line wraps to {n} lines: {line!r}"
    set_text(title, TITLE)

    set_text(find(slide, "TextBox 7"),
             ["Jiaming Wei          Supervisors: Prof. María Pérez-Ortiz  ·  Zekun Wu"])
    set_text(find(slide, "TextBox 8"),
             ["UCL Centre for Artificial Intelligence          Holistic AI"])
    standfirst = find(slide, "TextBox 12")
    n = line_count(STANDFIRST, SERIF, SZ_STANDFIRST.pt, int(standfirst.width * 0.94))
    assert n <= 2, f"standfirst wraps to {n} lines at 28pt and will leave its band"
    set_text(standfirst, [STANDFIRST])

    global HEADER_PARTS
    HEADER_PARTS = tuple(find(slide, n) for n in
                         ("Rectangle 13", "TextBox 14", "Rectangle 15", "Rectangle 16"))
    drop(slide, "TextBox 17", "Rectangle 33", "TextBox 34", "Rectangle 35",
         "Rectangle 36", "TextBox 37", "Rectangle 38", "Rectangle 39", "TextBox 42",
         "Rectangle 48", "TextBox 49", "Rectangle 50", "Rectangle 51",
         "Rectangle 52", "TextBox 53", "TextBox 54", "TextBox 55", "TextBox 56",
         "TextBox 57", "TextBox 58")

    loop_end = build_loop(slide)
    row_y = build_result_strip(slide, loop_end + Mm(9)) + Mm(9)
    ends = {"why (left)": (build_why(slide, row_y), ROW_BOTTOM),
            "evidence (right)": (build_evidence(slide, row_y), ROW_BOTTOM)}

    drop(slide, *HEADER_PARTS)

    set_text(find(slide, "TextBox 71"),
             ["jiaming.wei.22@ucl.ac.uk",
              "Paper, code & full results: github.com/Quarkgluonmixture/"
              "Cost-Aware-Routing-for-Web-Usage-Agents"])
    qr_png = FIGDIR / "qr_repo.png"
    qrcode.make(REPO_URL, box_size=20, border=1).save(qr_png)
    set_text(find(slide, "TextBox 73"), ["Explore all 8 settings →"])
    slot = find(slide, "Rectangle 74")
    slide.shapes.add_picture(str(qr_png), slot.left + Mm(2), slot.top + Mm(2),
                             width=slot.width - Mm(4))
    drop(slide, "Rectangle 74", "TextBox 75")

    prs.save(str(OUT))
    verify(prs, ends, row_y)


def verify(prs, ends, row_y):
    mm = lambda emu: emu / 914400 * 25.4  # noqa: E731
    w, h = mm(prs.slide_width), mm(prs.slide_height)
    assert abs(w - 594) < 0.5 and abs(h - 841) < 0.5, "slide was resized!"
    print(f"wrote {OUT.relative_to(REPO)}   ({w:.0f}x{h:.0f}mm, A1, not resized)")
    print(f"  {'row B starts':16s} {mm(row_y):6.1f}mm   (loop on top)")
    bad = False
    for name, (end, limit) in ends.items():
        slack = mm(limit - end)
        flag = "" if slack >= 0 else "   <-- OVERRUNS ITS BOX"
        bad |= slack < 0
        print(f"  {name:16s} ends {mm(end):6.1f}mm   (box to {mm(limit):.0f}mm, slack {slack:+6.1f}mm){flag}")
    if bad:
        raise SystemExit("a panel overran its box — shorten it or move the grid")


if __name__ == "__main__":
    main()
